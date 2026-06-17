from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import sys

# Content for the slides
slides_content = [
    {
        "title": "NeuroDraw",
        "subtitle": "Early Parkinson’s Disease Screening Using Spiral Drawing + AI",
        "content": [
            "Prepared by: Rahma Gamal Yousof – Miriam Bassem Boushra – Nouran Osama Ali",
            "New Mansoura University | Faculty of Computer Science and Engineering | 2025–2026"
        ],
        "layout": "Title Slide"
    },
    {
        "title": "Introduction",
        "content": [
            "Parkinson’s Disease (PD) is a progressive neurological disorder affecting motor control",
            "Early symptoms are subtle and often missed",
            "Early screening supports timely follow-up and better outcomes"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Problem Statement",
        "content": [
            "Diagnosis depends mainly on clinical observation (subjective)",
            "Limited access to specialists, especially in remote areas",
            "Need for an accessible, non-invasive early screening support tool"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Project Purpose",
        "content": [
            "Support early PD screening using spiral drawing analysis",
            "Provide a remote and low-cost decision-support solution",
            "Assist doctors in review and patient follow-up (not replace diagnosis)"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Objectives",
        "content": [
            "Train a CNN model to classify spiral drawings (PD vs Healthy)",
            "Achieve stable and reliable performance",
            "Integrate the model into a web-based platform",
            "Enable real-time doctor–patient communication and case tracking"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Related Work",
        "content": [
            "Classical ML: handcrafted features → limited generalization",
            "Deep learning (CNN): automatic feature learning from images",
            "Our focus: strong model + real system workflow integration"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Proposed System (NeuroDraw)",
        "content": [
            "Patient performs spiral test remotely",
            "ML model predicts screening outcome",
            "Doctor reviews results and manages follow-up",
            "Real-time chat supports continuous monitoring"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Dataset Overview",
        "content": [
            "Spiral drawing images",
            "Two classes: Parkinson’s / Healthy",
            "Balanced split: training and testing",
            "Binary image classification task"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Preprocessing",
        "content": [
            "Resize images to fixed CNN input size",
            "Normalize pixels to [0,1]",
            "Reduce noise/background artifacts for clearer patterns"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Data Augmentation",
        "content": [
            "Rotation, zooming, horizontal flipping",
            "Increases data diversity",
            "Reduces overfitting and improves generalization"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "CNN Architecture",
        "content": [
            "Convolution layers: extract tremor/irregularity patterns",
            "Pooling layers: reduce dimensions while preserving features",
            "Dense layers: final classification",
            "Sigmoid output: binary decision"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Training Strategy",
        "content": [
            "Trained over multiple epochs",
            "Monitored training/validation accuracy and loss",
            "Early stopping + regularization to prevent overfitting"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "System Architecture",
        "content": [
            "Front-End: patient & doctor dashboards, drawing/upload, results, chat",
            "Back-End: auth, storage, inference requests, notifications",
            "ML Service: spiral image inference endpoint",
            "Database: users, tests, results, chat history"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Experimental Results",
        "content": [
            "Training accuracy ≈ 95%",
            "Validation accuracy ≈ 87–88%",
            "Stable convergence and good generalization"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Evaluation Metrics",
        "content": [
            "Accuracy, Precision, Recall (Sensitivity), F1-score",
            "Confusion matrix to analyze FP/FN",
            "Medical screening priority: minimize false negatives"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Result Interpretation",
        "content": [
            "Spiral drawings reflect fine motor impairment",
            "CNN captures tremor and irregularity patterns",
            "Suitable for early screening decision-support"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Discussion",
        "content": [
            "Non-invasive and low-cost approach",
            "Reduces subjectivity in manual assessment",
            "Limitations: dataset size, device variability, need clinical validation"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Future Work (Smart Pen)",
        "content": [
            "Integrate smart pen/stylus signals: pressure, speed, tilt, acceleration",
            "Extract tremor frequency features (time-series analysis)",
            "Multi-modal AI: image + motion signals for higher robustness"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Future Work (Clinical + Deployment)",
        "content": [
            "Larger and more diverse datasets",
            "Multi-center clinical validation",
            "Explainable AI (Grad-CAM) to improve clinical trust",
            "Mobile/tablet support and hospital system integration (EHR)"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Conclusion",
        "content": [
            "NeuroDraw demonstrates feasibility of AI-based spiral screening",
            "Combines ML with a practical web platform",
            "Supports doctors in early screening and follow-up"
        ],
        "layout": "Bullet Points"
    },
    {
        "title": "Thank You",
        "content": [
            "Thank You for Your Attention"
        ],
        "layout": "Centered Text"
    }
]

def create_presentation():
    prs = Presentation()

    # Define layouts (0: Title, 1: Content)
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Just use Title Only and manual textbox for centered text for the last slide if needed
    # Or stick to Bullet Points and Title Slide

    for slide_data in slides_content:
        if slide_data["layout"] == "Title Slide":
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            if "subtitle" in slide_data:
                subtitle.text = slide_data["subtitle"] + "\n\n" + "\n".join(slide_data["content"])
            else:
                 subtitle.text = "\n".join(slide_data["content"])
        
        elif slide_data["layout"] == "Centered Text":
             # Use Title Only layout
             slide = prs.slides.add_slide(prs.slide_layouts[5]) # 5 is often Title Only or Blank
             # But if we use slide_layouts[1] it has a body
             # Let's simple reuse bullet layout but centering? No, let's keep it simple.
             slide = prs.slides.add_slide(bullet_slide_layout)
             title = slide.shapes.title
             title.text = slide_data["title"]
             body = slide.placeholders[1]
             tf = body.text_frame
             tf.text = "\n".join(slide_data["content"])
             # Align center
             for p in tf.paragraphs:
                 p.alignment = PP_ALIGN.CENTER
                 
        else: # Bullet Points
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            tf = body.text_frame
            tf.text = slide_data["content"][0] if slide_data["content"] else ""
            
            for line in slide_data["content"][1:]:
                p = tf.add_paragraph()
                p.text = line

    output_file = "NeuroDraw_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library is not installed.")
        print("Please run: pip install python-pptx")
    except Exception as e:
        print(f"An error occurred: {e}")
